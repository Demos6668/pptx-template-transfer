#!/usr/bin/env python3
"""PPTX Template Transfer — apply one deck's visual design to another's content.

Usage:
    python3 pptx_template_transfer.py template.pptx content.pptx output.pptx
    python3 pptx_template_transfer.py template.pptx content.pptx output.pptx --mode recreate
    python3 pptx_template_transfer.py template.pptx content.pptx output.pptx --verbose
    python3 pptx_template_transfer.py template.pptx content.pptx output.pptx --report report.json

Modes:
    recreate — Analyze template style, extract content, rebuild from scratch (default).
               Zero template text leakage, clean XML, works in any viewer.
    clone    — Clone template slides as visual skeletons, inject content text.
               (alias: design)
    layout   — Transfer theme + masters + layouts between files.
"""

from __future__ import annotations

import argparse
import io
import json
import logging
import math
import re
import sys
import traceback
import zipfile
from copy import deepcopy
from dataclasses import asdict, dataclass, field
from datetime import date
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Emu, Inches, Pt

log = logging.getLogger("pptx_template_transfer")

# ============================================================================
# CONFIGURATION
# ============================================================================

@dataclass(frozen=True)
class Thresholds:
    """All classification thresholds in one place — tune per template style."""
    # Shape role classification
    title_min_font_pt: float = 20.0
    title_max_words: int = 20
    body_min_area_pct: float = 4.0
    body_min_area_pct_relaxed: float = 3.0
    body_min_words: int = 10
    body_min_words_relaxed: int = 5
    body_max_zones: int = 2
    decorative_max_area_pct: float = 2.0
    decorative_max_words: int = 5
    decorative_max_font_pt: float = 10.0
    footer_bottom_frac: float = 0.90
    footer_top_frac: float = 0.08
    footer_max_area_pct: float = 5.0
    info_left_frac: float = 0.55
    info_min_words: int = 5
    info_max_words: int = 50
    # Content extraction
    image_min_area_pct: float = 3.0
    subheading_min_font_pt: float = 18.0
    # Matching
    variety_max_pct: float = 0.40
    # Overflow
    overflow_max_font_scale: float = 0.70
    overflow_chars_per_sq_inch: float = 180.0


@dataclass
class TransferConfig:
    mode: str | None = None
    verbose: bool = False
    slide_map: dict[str, int] | None = None
    preserve_notes: bool = True
    auto_split: bool = False
    thresholds: Thresholds = field(default_factory=Thresholds)
    report_path: Path | None = None


# ============================================================================
# XML HELPERS
# ============================================================================

_NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "dgm": "http://schemas.openxmlformats.org/drawingml/2006/diagram",
}

_FOOTER_PATTERNS = re.compile(
    r"(?i)(page\s*\d+|confidential|©|\bcopyright\b|\ball rights reserved\b"
    r"|\b\d{4}[-/]\d{2}[-/]\d{2}\b|\b\d{2}/\d{2}/\d{4}\b"
    r"|proprietary|internal use|draft|do not distribute)",
)
_PAGE_NUM_PATTERN = re.compile(r"(?i)page\s*\d+")
_DATE_PATTERN = re.compile(
    r"\b(\d{4}[-/]\d{1,2}[-/]\d{1,2}|\d{1,2}/\d{1,2}/\d{4})\b",
)


def _update_rids_in_tree(element, rid_map: dict[str, str]) -> None:
    for el in element.iter():
        for attr_name in list(el.attrib.keys()):
            val = el.attrib[attr_name]
            if val in rid_map:
                el.attrib[attr_name] = rid_map[val]


# ============================================================================
# SHAPE HELPERS
# ============================================================================

def _text_of(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def _word_count(text: str) -> int:
    return len(text.split()) if text else 0


def _max_font_pt(shape) -> float:
    mx = 0.0
    if not shape.has_text_frame:
        return mx
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None:
                mx = max(mx, run.font.size.pt)
    return mx


def _min_font_pt(shape) -> float:
    mn = 999.0
    if not shape.has_text_frame:
        return 0.0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.size is not None and run.font.size.pt > 0:
                mn = min(mn, run.font.size.pt)
    return mn if mn < 999.0 else 0.0


def _shape_area(shape) -> int:
    return (shape.width or 0) * (shape.height or 0)


def _shape_area_pct(shape, slide_w: int, slide_h: int) -> float:
    total = slide_w * slide_h
    return _shape_area(shape) / total * 100.0 if total else 0.0


def _is_picture(shape) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == 13
    except Exception:
        return False


def _is_table(shape) -> bool:
    return hasattr(shape, "has_table") and shape.has_table


def _is_chart(shape) -> bool:
    return hasattr(shape, "has_chart") and shape.has_chart


def _is_group(shape) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.GROUP
    except Exception:
        return False


def _is_ole_or_embedded(shape) -> bool:
    try:
        tag = etree.QName(shape._element.tag).localname
        return tag == "graphicFrame" and not _is_table(shape) and not _is_chart(shape)
    except Exception:
        return False


def _shape_bottom_frac(shape, slide_h: int) -> float:
    if slide_h == 0:
        return 0.0
    return ((shape.top or 0) + (shape.height or 0)) / slide_h


def _shape_top_frac(shape, slide_h: int) -> float:
    return (shape.top or 0) / slide_h if slide_h else 0.0


def _shape_left_frac(shape, slide_w: int) -> float:
    return (shape.left or 0) / slide_w if slide_w else 0.0


def _is_allcaps_short(text: str) -> bool:
    words = text.split()
    if not words or len(words) > 5:
        return False
    alpha = "".join(c for c in text if c.isalpha())
    return bool(alpha) and alpha == alpha.upper()


def _dominant_text_color(shape) -> str | None:
    if not shape.has_text_frame:
        return None
    colors: dict[str, int] = {}
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                c = run.font.color
                if c and c.type is not None and c.rgb:
                    key = str(c.rgb)
                    colors[key] = colors.get(key, 0) + len(run.text)
            except (AttributeError, TypeError):
                pass
    return max(colors, key=colors.get) if colors else None


def _group_text_words(shape) -> int:
    """Count total words in a group shape's child text frames."""
    total = 0
    try:
        for child in shape.shapes:
            total += _word_count(_text_of(child))
            if _is_group(child):
                total += _group_text_words(child)
    except Exception:
        pass
    return total


def _has_placeholder_type(shape, ph_types: set[int]) -> bool:
    """Check if shape is a placeholder with one of the given types."""
    try:
        pf = shape.placeholder_format
        if pf is not None and pf.type is not None:
            return int(pf.type) in ph_types
    except Exception:
        pass
    return False


def _placeholder_type_int(shape) -> int | None:
    try:
        pf = shape.placeholder_format
        if pf is not None and pf.type is not None:
            return int(pf.type)
    except Exception:
        pass
    return None


# Placeholder type constants (from pptx.enum.shapes.PP_PLACEHOLDER)
_PH_TITLE = 15       # PP_PLACEHOLDER.TITLE
_PH_CENTER_TITLE = 3 # PP_PLACEHOLDER.CENTER_TITLE
_PH_SUBTITLE = 4     # PP_PLACEHOLDER.SUBTITLE
_PH_BODY = 2         # PP_PLACEHOLDER.BODY
_PH_OBJECT = 7       # PP_PLACEHOLDER.OBJECT
_PH_FOOTER = 11      # PP_PLACEHOLDER.FOOTER
_PH_SLIDE_NUMBER = 12  # PP_PLACEHOLDER.SLIDE_NUMBER
_PH_DATE = 10        # PP_PLACEHOLDER.DATE

_PH_TITLE_SET = {_PH_TITLE, _PH_CENTER_TITLE}
_PH_BODY_SET = {_PH_BODY, _PH_OBJECT, _PH_SUBTITLE}
_PH_FOOTER_SET = {_PH_FOOTER, _PH_SLIDE_NUMBER, _PH_DATE}


# ============================================================================
# A. SHAPE ROLE CLASSIFICATION — with placeholder awareness + confidence
# ============================================================================

@dataclass
class ShapeInfo:
    """Pre-computed properties of a shape for classification."""
    shape: Any
    text: str
    word_count: int
    font_size: float
    area_pct: float
    top_frac: float
    bottom_frac: float
    left_frac: float
    is_picture: bool
    is_table: bool
    is_chart: bool
    is_group: bool
    is_ole: bool
    placeholder_type: int | None
    name_lower: str
    group_text_words: int = 0


def _precompute_shape_info(shape, slide_w: int, slide_h: int) -> ShapeInfo:
    text = _text_of(shape)
    is_grp = _is_group(shape)
    return ShapeInfo(
        shape=shape,
        text=text,
        word_count=_word_count(text),
        font_size=_max_font_pt(shape),
        area_pct=_shape_area_pct(shape, slide_w, slide_h),
        top_frac=_shape_top_frac(shape, slide_h),
        bottom_frac=_shape_bottom_frac(shape, slide_h),
        left_frac=_shape_left_frac(shape, slide_w),
        is_picture=_is_picture(shape),
        is_table=_is_table(shape),
        is_chart=_is_chart(shape),
        is_group=is_grp,
        is_ole=_is_ole_or_embedded(shape),
        placeholder_type=_placeholder_type_int(shape),
        name_lower=(shape.name or "").lower(),
        group_text_words=_group_text_words(shape) if is_grp else 0,
    )


def _classify_shape(
    si: ShapeInfo,
    th: Thresholds,
    *,
    largest_font: float,
    median_font: float,
    title_assigned: bool,
    body_count: int,
    info_count: int,
    similar_ids: set,
) -> tuple[str, float]:
    """Classify a shape's role. Returns (role, confidence)."""

    # --- Placeholder shortcut (highest confidence) ---
    if si.placeholder_type is not None:
        if si.placeholder_type in _PH_TITLE_SET:
            return ("title", 0.95) if not title_assigned else ("decorative", 0.6)
        if si.placeholder_type in _PH_BODY_SET:
            return ("body", 0.95) if body_count < th.body_max_zones else ("decorative", 0.5)
        if si.placeholder_type in _PH_FOOTER_SET:
            return ("footer", 0.95)

    # --- MEDIA ---
    if si.is_picture or si.is_chart or si.is_table:
        return ("media", 0.95)
    if si.is_ole:
        return ("media", 0.90)
    if si.is_group:
        # Group with substantial text → potential body zone
        if si.group_text_words > 20 and body_count < th.body_max_zones:
            return ("body", 0.55)
        return ("media", 0.80)

    # --- FOOTER / HEADER ---
    if si.bottom_frac >= th.footer_bottom_frac and si.area_pct < th.footer_max_area_pct:
        return ("footer", 0.85)
    if si.top_frac <= th.footer_top_frac and si.area_pct < 3 and si.word_count <= 10:
        return ("footer", 0.80)
    if si.text and _FOOTER_PATTERNS.search(si.text):
        return ("footer", 0.85)

    # --- NO TEXT → DECORATIVE ---
    if not si.text:
        return ("decorative", 0.90)

    # --- DECORATIVE checks ---
    if si.area_pct < th.decorative_max_area_pct and si.word_count <= th.decorative_max_words:
        return ("decorative", 0.80)
    if 0 < si.font_size <= th.decorative_max_font_pt:
        return ("decorative", 0.75)
    if si.word_count <= 3:
        return ("decorative", 0.70)
    if _is_allcaps_short(si.text) and si.area_pct < 5:
        return ("decorative", 0.75)
    if re.match(r"^\d{1,2}$", si.text.strip()):
        return ("decorative", 0.90)
    if id(si.shape) in similar_ids:
        return ("decorative", 0.70)

    # --- Adaptive title threshold ---
    # Use 80th percentile if median font is low (dense slide)
    effective_title_font = th.title_min_font_pt
    if median_font > 0 and median_font < 14:
        effective_title_font = max(median_font * 1.3, 14)

    # --- TITLE ---
    conf = 0.0
    if not title_assigned and si.top_frac < 0.45 and si.word_count <= th.title_max_words:
        if si.font_size >= effective_title_font and si.font_size >= largest_font - 2:
            conf = 0.85
        # Name boost
        if any(kw in si.name_lower for kw in ("title", "heading")):
            conf = max(conf, 0.70)
        if conf >= 0.55:
            return ("title", conf)

    # --- BODY ---
    if body_count < th.body_max_zones:
        if si.area_pct > th.body_min_area_pct and si.word_count > th.body_min_words:
            conf = 0.80
            if any(kw in si.name_lower for kw in ("body", "content", "text")):
                conf = 0.85
            return ("body", conf)
        if si.area_pct > th.body_min_area_pct_relaxed and si.word_count > th.body_min_words_relaxed:
            return ("body", 0.60)

    # --- INFO (sidebar/panel) ---
    if (info_count < 1
            and si.left_frac >= th.info_left_frac
            and th.info_min_words <= si.word_count <= th.info_max_words
            and si.area_pct > 2):
        return ("info", 0.65)

    return ("decorative", 0.40)


def _detect_repeated_patterns(infos: list[ShapeInfo], slide_w: int, slide_h: int) -> set:
    """Find shape ids that are part of repeated visual patterns."""
    result: set[int] = set()
    if len(infos) < 3:
        return result

    dimension_groups: dict[tuple, list] = {}
    for si in infos:
        w = si.shape.width or 0
        h = si.shape.height or 0
        if w == 0 or h == 0:
            continue
        bw = round(w / (slide_w * 0.02)) if slide_w else 0
        bh = round(h / (slide_h * 0.02)) if slide_h else 0
        dimension_groups.setdefault((bw, bh), []).append(si)

    for group in dimension_groups.values():
        if len(group) < 3:
            continue
        top_buckets: dict[int, int] = {}
        for si in group:
            bucket = round((si.shape.top or 0) / (slide_h * 0.05)) if slide_h else 0
            top_buckets[bucket] = top_buckets.get(bucket, 0) + 1
        if top_buckets and max(top_buckets.values()) >= 3:
            for si in group:
                if si.word_count <= 15:
                    result.add(id(si.shape))
    return result


def classify_all_shapes(
    slide, slide_w: int, slide_h: int, th: Thresholds,
) -> list[tuple[Any, str, float]]:
    """Classify all shapes. Returns [(shape, role, confidence), ...]."""
    shapes = list(slide.shapes)

    # Pre-compute all shape info in one pass
    infos = [_precompute_shape_info(s, slide_w, slide_h) for s in shapes]

    # Slide-level stats
    fonts = [si.font_size for si in infos if si.font_size > 0]
    largest_font = max(fonts) if fonts else 0.0
    sorted_fonts = sorted(fonts)
    median_font = sorted_fonts[len(sorted_fonts) // 2] if sorted_fonts else 0.0

    similar_ids = _detect_repeated_patterns(infos, slide_w, slide_h)

    # Sort by position for classification priority
    sorted_infos = sorted(infos, key=lambda si: ((si.shape.top or 0), (si.shape.left or 0)))

    title_assigned = False
    body_count = 0
    info_count = 0
    results: dict[int, tuple[str, float]] = {}

    for si in sorted_infos:
        role, conf = _classify_shape(
            si, th,
            largest_font=largest_font,
            median_font=median_font,
            title_assigned=title_assigned,
            body_count=body_count,
            info_count=info_count,
            similar_ids=similar_ids,
        )
        results[id(si.shape)] = (role, conf)
        if role == "title":
            title_assigned = True
        elif role == "body":
            body_count += 1
        elif role == "info":
            info_count += 1

    return [(s, *results[id(s)]) for s in shapes]


def classify_shape_role(
    shape, slide_width: int, slide_height: int,
    slide=None, th: Thresholds | None = None,
) -> str:
    """Public single-shape classifier. Returns "title"|"body"|"decorative"|"footer".

    Note: For accurate title/body assignment (which depends on slide context),
    prefer classify_all_shapes() or get_slide_zones() instead.
    """
    if th is None:
        th = Thresholds()
    si = _precompute_shape_info(shape, slide_width, slide_height)
    role, _ = _classify_shape(
        si, th, largest_font=si.font_size, median_font=si.font_size,
        title_assigned=False, body_count=0, info_count=0, similar_ids=set(),
    )
    # Collapse "media" and "info" into the 4-role scheme
    if role == "media":
        return "decorative"
    if role == "info":
        return "body"
    return role


def get_slide_zones(
    slide, slide_width: int, slide_height: int, th: Thresholds | None = None,
) -> dict[str, list]:
    """Classify every shape on a slide and return grouped zones.

    Returns {"title": [shapes], "body": [shapes], "decorative": [shapes], "footer": [shapes]}
    """
    if th is None:
        th = Thresholds()
    classifications = classify_all_shapes(slide, slide_width, slide_height, th)
    zones: dict[str, list] = {"title": [], "body": [], "decorative": [], "footer": []}
    for shape, role, _conf in classifications:
        if role == "title":
            zones["title"].append(shape)
        elif role in ("body", "info"):
            zones["body"].append(shape)
        elif role == "footer":
            zones["footer"].append(shape)
        else:  # decorative, media
            zones["decorative"].append(shape)
    return zones


# ============================================================================
# B. CONTENT STRUCTURE EXTRACTOR
# ============================================================================

@dataclass
class RunData:
    text: str
    bold: bool = False
    italic: bool = False
    font_size: float = 0.0
    hyperlink_url: str | None = None


@dataclass
class ParagraphData:
    text: str
    level: int = 0
    bold: bool = False
    italic: bool = False
    font_size: float = 0.0
    runs: list[RunData] = field(default_factory=list)


@dataclass
class TextBlock:
    """A positioned text group preserving spatial layout from the source slide."""
    paragraphs: list[ParagraphData] = field(default_factory=list)
    left_pct: float = 0.0   # % of slide width
    top_pct: float = 0.0    # % of slide height
    width_pct: float = 20.0
    height_pct: float = 10.0
    is_heading: bool = False  # bold/large short text
    is_label: bool = False    # very short text (≤3 words), e.g. numbers, tags


@dataclass
class ContentData:
    title: str = ""
    body_paragraphs: list[ParagraphData] = field(default_factory=list)
    text_blocks: list[TextBlock] = field(default_factory=list)
    tables: list[dict] = field(default_factory=list)
    images: list[tuple] = field(default_factory=list)
    charts: list[Any] = field(default_factory=list)  # chart element + rels
    has_chart: bool = False
    slide_type: str = "content"
    word_count: int = 0
    primary_color: str | None = None
    notes: str = ""


def _classify_slide_type(slide, slide_index: int, total_slides: int) -> str:
    texts = []
    images = tables = charts = 0
    for shape in slide.shapes:
        if _is_picture(shape):
            images += 1
        if _is_table(shape):
            tables += 1
        if _is_chart(shape):
            charts += 1
        t = _text_of(shape)
        if t:
            texts.append({"size": _max_font_pt(shape), "words": _word_count(t)})

    total_words = sum(t["words"] for t in texts)
    big = [t for t in texts if t["size"] >= 20]

    if not texts and images == 0:
        return "blank"
    if not texts and images > 0:
        return "image"
    if slide_index == 0 and big:
        return "title"
    if total_words <= 20 and big and len(texts) <= 5:
        return "title"
    if slide_index == total_slides - 1 and total_words <= 40:
        return "closing"
    if len(texts) <= 3 and total_words <= 15 and big:
        return "section"
    if tables > 0 or charts > 0:
        return "data"
    if images >= 3 and total_words < 30:
        return "image"
    return "content"


def _extract_paragraphs_from_shape(shape) -> list[ParagraphData]:
    result = []
    if not shape.has_text_frame:
        return result
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        level = para.level if para.level else 0
        bold = italic = False
        font_size = 0.0
        runs_data = []
        for run in para.runs:
            r_bold = bool(run.font.bold)
            r_italic = bool(run.font.italic)
            r_size = run.font.size.pt if run.font.size else 0.0
            r_url = None
            try:
                if run.hyperlink and run.hyperlink.address:
                    r_url = run.hyperlink.address
            except Exception:
                pass
            runs_data.append(RunData(
                text=run.text, bold=r_bold, italic=r_italic,
                font_size=r_size, hyperlink_url=r_url,
            ))
            if r_bold:
                bold = True
            if r_italic:
                italic = True
            font_size = max(font_size, r_size)
        result.append(ParagraphData(
            text=text, level=level, bold=bold,
            italic=italic, font_size=font_size, runs=runs_data,
        ))
    return result


def _extract_table_data(shape) -> list[list[str]]:
    if not _is_table(shape):
        return []
    rows = []
    for row in shape.table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return rows


def _extract_chart_info(shape, slide) -> dict | None:
    """Extract chart element and relationship info for cloning."""
    if not _is_chart(shape):
        return None
    try:
        chart_part = shape.chart_part
        return {
            "element": deepcopy(shape._element),
            "chart_part": chart_part,
            "width": shape.width,
            "height": shape.height,
            "left": shape.left,
            "top": shape.top,
        }
    except Exception:
        return None


def extract_content(
    slide, slide_index: int, total_slides: int,
    slide_w: int, slide_h: int, th: Thresholds,
) -> ContentData:
    content = ContentData()
    content.slide_type = _classify_slide_type(slide, slide_index, total_slides)
    shapes = list(slide.shapes)

    # --- Title detection ---
    text_shapes = [(s, _max_font_pt(s), _text_of(s)) for s in shapes if _text_of(s)]
    text_shapes.sort(key=lambda x: (-x[1], (x[0].top or 0)))

    title_shape = None
    for s, fs, txt in text_shapes:
        if _word_count(txt) <= 15 and fs >= 20:
            title_shape = s
            content.title = txt
            break
    if not title_shape and text_shapes:
        for s, fs, txt in sorted(text_shapes, key=lambda x: (x[0].top or 0)):
            if _word_count(txt) <= 10:
                title_shape = s
                content.title = txt
                break

    # --- Body extraction (flat paragraphs + positioned text blocks) ---
    body_shapes = [
        s for s in shapes
        if s is not title_shape and _text_of(s) and not _is_table(s) and not _is_chart(s)
    ]
    # Filter out footer-zone shapes from body
    body_shapes = [
        s for s in body_shapes
        if _shape_bottom_frac(s, slide_h) <= 0.92
        and not _FOOTER_PATTERNS.match(_text_of(s).strip())
        and not (_placeholder_type_int(s) is not None and _placeholder_type_int(s) in _PH_FOOTER_SET)
    ]
    body_shapes.sort(key=lambda s: ((s.top or 0), (s.left or 0)))

    for shape in body_shapes:
        paras = _extract_paragraphs_from_shape(shape)
        for p in paras:
            if p.bold or (p.font_size >= th.subheading_min_font_pt and _word_count(p.text) <= 10):
                p.bold = True
            content.body_paragraphs.append(p)

        # Also build positioned TextBlock for recreate mode
        if paras and slide_w > 0 and slide_h > 0:
            text = _text_of(shape)
            wc = _word_count(text)
            max_fs = _max_font_pt(shape)
            is_heading = (max_fs >= th.subheading_min_font_pt and wc <= 10) or (
                paras[0].bold and wc <= 10
            )
            is_label = wc <= 3 and max_fs < 20
            content.text_blocks.append(TextBlock(
                paragraphs=paras,
                left_pct=(shape.left or 0) / slide_w * 100,
                top_pct=(shape.top or 0) / slide_h * 100,
                width_pct=(shape.width or 0) / slide_w * 100,
                height_pct=(shape.height or 0) / slide_h * 100,
                is_heading=is_heading,
                is_label=is_label,
            ))

    # --- Tables ---
    for shape in shapes:
        if _is_table(shape):
            table_text = _extract_table_data(shape)
            content.tables.append({
                "data": table_text,
                "rows": len(table_text),
                "cols": len(table_text[0]) if table_text else 0,
                "element": deepcopy(shape._element),
                "width": shape.width, "height": shape.height,
                "left": shape.left, "top": shape.top,
            })

    # --- Charts ---
    for shape in shapes:
        if _is_chart(shape):
            content.has_chart = True
            ci = _extract_chart_info(shape, slide)
            if ci:
                content.charts.append(ci)

    # --- Images ---
    for shape in shapes:
        if _is_picture(shape):
            area_pct = _shape_area_pct(shape, slide_w, slide_h)
            if area_pct > th.image_min_area_pct:
                try:
                    blob = shape.image.blob
                    content.images.append(
                        (blob, shape.width, shape.height, shape.left or 0, shape.top or 0),
                    )
                except Exception:
                    pass

    # --- Speaker notes ---
    try:
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                content.notes = notes_tf.text.strip()
    except Exception:
        pass

    # --- Word count and color ---
    all_text = content.title + " " + " ".join(p.text for p in content.body_paragraphs)
    content.word_count = _word_count(all_text)
    colors: dict[str, int] = {}
    for shape in shapes:
        c = _dominant_text_color(shape)
        if c:
            colors[c] = colors.get(c, 0) + 1
    if colors:
        content.primary_color = max(colors, key=colors.get)

    return content


# ============================================================================
# C. SMART SLIDE MATCHING
# ============================================================================

def _classify_template_structure(
    slide, slide_w: int, slide_h: int,
    slide_index: int = -1, total_slides: int = -1,
) -> str:
    shapes = list(slide.shapes)
    text_shapes = [s for s in shapes if _text_of(s)]
    images = [s for s in shapes if _is_picture(s)]
    tables = [s for s in shapes if _is_table(s)]
    charts = [s for s in shapes if _is_chart(s)]

    total_words = sum(_word_count(_text_of(s)) for s in text_shapes)
    big = [s for s in text_shapes if _max_font_pt(s) >= 20]

    if tables or charts:
        return "data"
    if len(images) >= 2 and total_words < 30:
        return "visual"

    numbered = [s for s in text_shapes if re.match(r"^\d{1,2}$", _text_of(s).strip())]
    if len(numbered) >= 3:
        infos = [_precompute_shape_info(s, slide_w, slide_h) for s in shapes]
        if _detect_repeated_patterns(infos, slide_w, slide_h):
            return "grid"
        return "list"

    if slide_index == 0 and big:
        return "title"

    closing_kw = {"thank", "contact", "questions", "q&a", "reference"}
    all_lower = " ".join(_text_of(s) for s in text_shapes).lower()
    if slide_index == total_slides - 1 and total_slides > 1:
        if any(kw in all_lower for kw in closing_kw) or total_words <= 40:
            return "closing"

    if big and total_words <= 20 and len(text_shapes) <= 5:
        return "title"
    if big and len(text_shapes) <= 3 and total_words <= 15:
        return "section"
    if any(kw in all_lower for kw in closing_kw) and total_words <= 40:
        return "closing"

    body = [s for s in text_shapes if _shape_area_pct(s, slide_w, slide_h) > 4 and _word_count(_text_of(s)) > 10]
    if body:
        return "narrative"
    return "narrative" if total_words > 20 else "section"


_TYPE_COMPAT = {
    ("title", "title"): 40, ("title", "section"): 20, ("title", "narrative"): 10,
    ("content", "narrative"): 40, ("content", "list"): 30, ("content", "grid"): 25,
    ("content", "data"): 15, ("content", "title"): 8,
    ("section", "section"): 40, ("section", "title"): 25,
    ("data", "data"): 40, ("data", "narrative"): 20, ("data", "grid"): 25,
    ("closing", "closing"): 40, ("closing", "section"): 20,
    ("image", "visual"): 40, ("image", "narrative"): 15,
    ("blank", "section"): 10,
}


def _match_score(
    c_type: str, t_struct: str,
    ci: int, ti: int, ct: int, tt: int,
    c_words: int, t_words: int,
    c_has_table: bool, t_has_table: bool,
    c_paras: int, t_is_list: bool,
) -> float:
    score = float(_TYPE_COMPAT.get((c_type, t_struct), 5))

    if t_words > 0 and c_words > 0:
        score += 25 * min(c_words, t_words) / max(c_words, t_words)
    elif c_words == 0 and t_words <= 10:
        score += 20

    if c_has_table and t_has_table:
        score += 20
    elif c_has_table:
        score += 5
    elif c_paras >= 5 and t_is_list:
        score += 15
    elif c_paras >= 3 and t_struct == "narrative":
        score += 15
    elif c_paras < 3 and t_struct in ("section", "title"):
        score += 10

    if ci == 0 and t_struct == "title":
        score += 15
    elif ci == ct - 1 and t_struct == "closing":
        score += 15
    elif ct > 1 and tt > 1:
        score += 15 * (1 - abs(ci / (ct - 1) - ti / (tt - 1)))

    return score


def build_slide_mapping(
    content_prs: Presentation, template_prs: Presentation,
    content_data_list: list[ContentData], th: Thresholds,
) -> list[int]:
    sw, sh = template_prs.slide_width, template_prs.slide_height
    ct, tt = len(content_prs.slides), len(template_prs.slides)

    t_info = []
    for i, slide in enumerate(template_prs.slides):
        struct = _classify_template_structure(slide, sw, sh, i, tt)
        words = sum(_word_count(_text_of(s)) for s in slide.shapes)
        t_info.append({
            "struct": struct, "words": words,
            "has_table": any(_is_table(s) for s in slide.shapes),
            "is_list": struct in ("list", "grid"),
        })

    # Score matrix
    score_matrix: list[list[tuple[int, float]]] = []
    for ci, cd in enumerate(content_data_list):
        scores = []
        for ti, tinfo in enumerate(t_info):
            sc = _match_score(
                cd.slide_type, tinfo["struct"], ci, ti, ct, tt,
                cd.word_count, tinfo["words"],
                len(cd.tables) > 0, tinfo["has_table"],
                len(cd.body_paragraphs), tinfo["is_list"],
            )
            scores.append((ti, sc))
        scores.sort(key=lambda x: -x[1])
        score_matrix.append(scores)

    # Greedy with variety
    usage: dict[int, int] = {i: 0 for i in range(tt)}
    max_per = max(2, math.ceil(ct * th.variety_max_pct))
    min_distinct = min(tt, max(3, math.ceil(ct / 3)))

    mapping = []
    for ci, scores in enumerate(score_matrix):
        best_idx, best_sc = scores[0]
        if usage[best_idx] >= max_per:
            for ti2, sc2 in scores[1:]:
                if usage[ti2] < max_per:
                    best_idx = ti2
                    break
        mapping.append(best_idx)
        usage[best_idx] += 1

    # Redistribute to hit min_distinct
    used_set = {ti for ti, c in usage.items() if c > 0}
    unused = [ti for ti in range(tt) if usage[ti] == 0]
    if len(used_set) < min_distinct and unused:
        overused = sorted(
            [(ti, c) for ti, c in usage.items() if c > 1], key=lambda x: -x[1],
        )
        for u_ti in unused:
            if not overused:
                break
            donor_ti = overused[0][0]
            candidates = [(ci, score_matrix[ci]) for ci in range(ct) if mapping[ci] == donor_ti]
            best_ci, best_sc = None, -1.0
            for ci2, scores in candidates:
                for ti2, sc2 in scores:
                    if ti2 == u_ti and sc2 > best_sc:
                        best_ci, best_sc = ci2, sc2
            if best_ci is not None and best_sc > 10:
                usage[mapping[best_ci]] -= 1
                mapping[best_ci] = u_ti
                usage[u_ti] = usage.get(u_ti, 0) + 1
                overused = sorted(
                    [(ti, c) for ti, c in usage.items() if c > 1], key=lambda x: -x[1],
                )

    for ci in range(ct):
        sc_val = next((sc for t, sc in score_matrix[ci] if t == mapping[ci]), 0)
        log.debug(
            "  Slide %d (%s, %dw) -> Template %d (%s) score=%.0f",
            ci + 1, content_data_list[ci].slide_type,
            content_data_list[ci].word_count, mapping[ci] + 1,
            t_info[mapping[ci]]["struct"], sc_val,
        )
    used = sum(1 for v in usage.values() if v > 0)
    log.debug("  Variety: %d/%d templates used (target: %d+)", used, tt, min_distinct)

    return mapping


# ============================================================================
# SLIDE CLONING
# ============================================================================

def _clone_slide(template_prs: Presentation, src_slide, dst_prs: Presentation):
    dst_layout = dst_prs.slide_layouts[0]
    src_layout_name = src_slide.slide_layout.name
    for layout in dst_prs.slide_layouts:
        if layout.name == src_layout_name:
            dst_layout = layout
            break

    new_slide = dst_prs.slides.add_slide(dst_layout)

    spTree = new_slide.shapes._spTree
    for sp in list(spTree):
        tag = etree.QName(sp.tag).localname if isinstance(sp.tag, str) else ""
        if tag in ("sp", "pic", "grpSp", "cxnSp", "graphicFrame"):
            spTree.remove(sp)

    src_spTree = src_slide.shapes._spTree
    for child in src_spTree:
        tag = etree.QName(child.tag).localname if isinstance(child.tag, str) else ""
        if tag in ("sp", "pic", "grpSp", "cxnSp", "graphicFrame"):
            spTree.append(deepcopy(child))

    # Background
    src_sld, dst_sld = src_slide._element, new_slide._element
    src_bg = src_sld.find(f'{{{_NSMAP["p"]}}}bg')
    if src_bg is not None:
        dst_bg = dst_sld.find(f'{{{_NSMAP["p"]}}}bg')
        if dst_bg is not None:
            dst_sld.remove(dst_bg)
        new_bg = deepcopy(src_bg)
        cSld = dst_sld.find(f'{{{_NSMAP["p"]}}}cSld')
        if cSld is not None:
            dst_sld.insert(list(dst_sld).index(cSld), new_bg)
        else:
            dst_sld.insert(0, new_bg)

    # Transition
    ns_p = _NSMAP["p"]
    src_transition = src_sld.find(f'{{{ns_p}}}transition')
    if src_transition is not None:
        dst_sld.append(deepcopy(src_transition))

    # Relationships
    rid_map: dict[str, str] = {}
    broken_rels: list[str] = []
    for rel_key, rel in src_slide.part.rels.items():
        if rel.reltype in (RT.SLIDE_LAYOUT, RT.NOTES_SLIDE):
            continue
        try:
            if rel.is_external:
                new_rid = new_slide.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
            else:
                new_rid = new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
            rid_map[rel_key] = new_rid
        except Exception as exc:
            broken_rels.append(f"{rel_key} ({rel.reltype}): {exc}")

    if broken_rels:
        log.warning("  Broken relationships during clone: %s", "; ".join(broken_rels))

    if rid_map:
        _update_rids_in_tree(spTree, rid_map)
        dst_bg2 = dst_sld.find(f'{{{_NSMAP["p"]}}}bg')
        if dst_bg2 is not None:
            _update_rids_in_tree(dst_bg2, rid_map)

    return new_slide


# ============================================================================
# D. TEXT INJECTION — multi-level format preservation
# ============================================================================

def _save_all_paragraph_formats(shape) -> dict[int, tuple]:
    """Save formatting from all paragraph levels in a shape.

    Returns dict mapping indent level -> (pPr_element, rPr_element).
    """
    ns_a = _NSMAP["a"]
    formats: dict[int, tuple] = {}
    if not shape.has_text_frame:
        return formats

    for para in shape.text_frame.paragraphs:
        p_el = para._p
        level = para.level if para.level else 0
        if level in formats:
            continue

        pPr = p_el.find(f'{{{ns_a}}}pPr')
        pPr_copy = deepcopy(pPr) if pPr is not None else None

        rPr_copy = None
        for r in p_el.findall(f'{{{ns_a}}}r'):
            rPr = r.find(f'{{{ns_a}}}rPr')
            if rPr is not None:
                rPr_copy = deepcopy(rPr)
                break
        if rPr_copy is None:
            endRPr = p_el.find(f'{{{ns_a}}}endParaRPr')
            if endRPr is not None:
                rPr_copy = deepcopy(endRPr)

        formats[level] = (pPr_copy, rPr_copy)

    return formats


def _get_format_for_level(formats: dict[int, tuple], level: int) -> tuple:
    """Get (pPr, rPr) for a given indent level, falling back to closest."""
    if level in formats:
        return formats[level]
    if not formats:
        return (None, None)
    closest = min(formats.keys(), key=lambda k: abs(k - level))
    return formats[closest]


def _estimate_text_capacity(shape, slide_w: int, slide_h: int, th: Thresholds) -> int:
    """Estimate how many characters a text shape can hold."""
    w_inches = (shape.width or 0) / 914400.0
    h_inches = (shape.height or 0) / 914400.0
    area_sq_inches = w_inches * h_inches
    if area_sq_inches <= 0:
        return 100
    return max(20, int(area_sq_inches * th.overflow_chars_per_sq_inch))


def _fit_paragraphs(
    paragraphs: list[ParagraphData], max_chars: int,
) -> list[ParagraphData]:
    """Truncate paragraphs to fit within max_chars, adding '...' if needed."""
    result = []
    chars_used = 0
    for p in paragraphs:
        if chars_used + len(p.text) <= max_chars:
            result.append(p)
            chars_used += len(p.text)
        else:
            remaining = max_chars - chars_used
            if remaining > 20:
                truncated = ParagraphData(
                    text=p.text[:remaining - 3] + "...",
                    level=p.level, bold=p.bold, italic=p.italic,
                    font_size=p.font_size, runs=p.runs,
                )
                result.append(truncated)
            elif not result:
                # At least include one truncated paragraph
                result.append(ParagraphData(
                    text=p.text[:max(50, max_chars)] + "...",
                    level=p.level, bold=p.bold,
                ))
            break
    return result


def _inject_text_simple(shape, text: str) -> None:
    """Replace text preserving first paragraph's formatting."""
    if not shape.has_text_frame or not shape.text_frame.paragraphs:
        return
    ns_a = _NSMAP["a"]
    formats = _save_all_paragraph_formats(shape)
    pPr, rPr = _get_format_for_level(formats, 0)

    txBody = shape.text_frame._txBody
    for p in list(txBody.findall(f'{{{ns_a}}}p')):
        txBody.remove(p)

    for para_text in text.split("\n"):
        new_p = etree.SubElement(txBody, f'{{{ns_a}}}p')
        if pPr is not None:
            new_p.append(deepcopy(pPr))
        if para_text.strip():
            new_r = etree.SubElement(new_p, f'{{{ns_a}}}r')
            if rPr is not None:
                new_r.append(deepcopy(rPr))
            new_t = etree.SubElement(new_r, f'{{{ns_a}}}t')
            new_t.text = para_text
        else:
            eRPr = etree.SubElement(new_p, f'{{{ns_a}}}endParaRPr')
            if rPr is not None:
                for k, v in rPr.attrib.items():
                    eRPr.attrib[k] = v


def _inject_structured_text(
    shape, paragraphs: list[ParagraphData], th: Thresholds,
    slide_w: int = 0, slide_h: int = 0,
) -> None:
    """Inject structured paragraphs with multi-level format preservation."""
    if not shape.has_text_frame or not paragraphs:
        return
    if not shape.text_frame.paragraphs:
        return

    ns_a = _NSMAP["a"]

    # Save formatting per level
    formats = _save_all_paragraph_formats(shape)

    # Overflow prevention
    if slide_w > 0 and slide_h > 0:
        capacity = _estimate_text_capacity(shape, slide_w, slide_h, th)
        total_chars = sum(len(p.text) for p in paragraphs)
        if total_chars > capacity:
            paragraphs = _fit_paragraphs(paragraphs, capacity)
            log.debug("    Overflow: truncated %d -> %d chars (capacity=%d)",
                       total_chars, sum(len(p.text) for p in paragraphs), capacity)

    txBody = shape.text_frame._txBody
    for p in list(txBody.findall(f'{{{ns_a}}}p')):
        txBody.remove(p)

    for pd in paragraphs:
        new_p = etree.SubElement(txBody, f'{{{ns_a}}}p')
        pPr_tmpl, rPr_tmpl = _get_format_for_level(formats, pd.level)

        if pPr_tmpl is not None:
            pPr = deepcopy(pPr_tmpl)
            if pd.level > 0:
                pPr.set("lvl", str(pd.level))
            new_p.append(pPr)
        elif pd.level > 0:
            pPr = etree.SubElement(new_p, f'{{{ns_a}}}pPr')
            pPr.set("lvl", str(pd.level))

        if pd.text.strip():
            # If we have runs with hyperlinks, use multi-run injection
            if pd.runs and any(r.hyperlink_url for r in pd.runs):
                for rd in pd.runs:
                    if not rd.text:
                        continue
                    new_r = etree.SubElement(new_p, f'{{{ns_a}}}r')
                    if rPr_tmpl is not None:
                        rPr = deepcopy(rPr_tmpl)
                        if pd.bold or rd.bold:
                            rPr.set("b", "1")
                        if pd.italic or rd.italic:
                            rPr.set("i", "1")
                        # Hyperlink
                        if rd.hyperlink_url:
                            hlinkClick = etree.SubElement(
                                rPr, f'{{{ns_a}}}hlinkClick',
                            )
                            hlinkClick.set(
                                f'{{{_NSMAP["r"]}}}id', "",  # Will need rel
                            )
                        new_r.append(rPr)
                    new_t = etree.SubElement(new_r, f'{{{ns_a}}}t')
                    new_t.text = rd.text
            else:
                new_r = etree.SubElement(new_p, f'{{{ns_a}}}r')
                if rPr_tmpl is not None:
                    rPr = deepcopy(rPr_tmpl)
                    if pd.bold:
                        rPr.set("b", "1")
                    if pd.italic:
                        rPr.set("i", "1")
                    new_r.append(rPr)
                new_t = etree.SubElement(new_r, f'{{{ns_a}}}t')
                new_t.text = pd.text
        else:
            eRPr = etree.SubElement(new_p, f'{{{ns_a}}}endParaRPr')
            if rPr_tmpl is not None:
                for k, v in rPr_tmpl.attrib.items():
                    eRPr.attrib[k] = v


def _clear_shape_text(shape) -> None:
    if not shape.has_text_frame:
        return
    ns_a = _NSMAP["a"]
    txBody = shape.text_frame._txBody
    for p in list(txBody.findall(f'{{{ns_a}}}p')):
        txBody.remove(p)
    new_p = etree.SubElement(txBody, f'{{{ns_a}}}p')
    etree.SubElement(new_p, f'{{{ns_a}}}endParaRPr')


# ---- Template text clearing: is_protected + prepare_cloned_slide ----

_JUST_NUMBER_RE = re.compile(r"^\d{1,2}$")


def _is_protected_shape(shape, slide_w: int, slide_h: int) -> bool:
    """Return True if the shape should keep its template text untouched.

    Aggressively clears ALL template text except truly structural elements
    (footers, page numbers, dates, confidential notices, empty shapes).
    This ensures zero template words leak into the output.
    """
    # No text frame — nothing to clear
    if not shape.has_text_frame:
        return True

    # Media shapes (picture, chart, table, group, OLE)
    if _is_picture(shape) or _is_chart(shape) or _is_table(shape) or _is_group(shape):
        return True
    if _is_ole_or_embedded(shape):
        return True

    # Empty text — nothing to clear
    text = shape.text_frame.text.strip()
    if not text:
        return True

    # Placeholder-based footer (slide number, date, footer) — structural
    ph = _placeholder_type_int(shape)
    if ph is not None and ph in _PH_FOOTER_SET:
        return True

    # Footer zone (bottom 8% of slide) — structural
    bottom_frac = _shape_bottom_frac(shape, slide_h)
    if bottom_frac > 0.92:
        return True

    # Common footer/label patterns (Page XX, Confidential, dates, ©)
    if _FOOTER_PATTERNS.match(text.strip()):
        return True

    # Just a number like "01", "02"
    if _JUST_NUMBER_RE.match(text.strip()):
        return True

    # NOT protected — this is an injection target whose text gets erased
    return False


def _prepare_cloned_slide(
    slide, slide_w: int, slide_h: int,
) -> tuple[list, list]:
    """Erase template text from injection targets, leave protected shapes untouched.

    Returns (injection_targets, protected_shapes) for diagnostic tracking.
    """
    targets: list = []
    protected: list = []

    for shape in slide.shapes:
        if _is_protected_shape(shape, slide_w, slide_h):
            protected.append(shape)
            continue

        # This is an injection target — CLEAR its text
        targets.append(shape)
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = ""

    return targets, protected


def inject_content(
    cloned_slide, content_data: ContentData,
    slide_w: int, slide_h: int, th: Thresholds,
) -> dict[str, Any]:
    """Inject content into a cloned slide. Returns diagnostic dict."""
    diag: dict[str, Any] = {
        "shapes": [], "injected_title": None,
        "injected_body": None, "protected_count": 0,
        "cleared_count": 0,
    }

    # Step 1: Classify shapes BEFORE clearing (needs original text for accuracy)
    classifications = classify_all_shapes(cloned_slide, slide_w, slide_h, th)

    # Step 2: Determine which shapes are injection targets vs protected
    targets, protected = _prepare_cloned_slide(cloned_slide, slide_w, slide_h)
    target_ids = {id(s) for s in targets}
    diag["cleared_count"] = len(targets)
    diag["protected_count"] = len(protected)

    # Step 3: Use pre-clearing classifications to assign title/body zones
    title_shape = None
    body_shapes: list = []

    for shape, role, conf in classifications:
        diag["shapes"].append({
            "name": shape.name, "role": role,
            "confidence": round(conf, 2),
            "area_pct": round(_shape_area_pct(shape, slide_w, slide_h), 1),
            "top_pct": round(_shape_top_frac(shape, slide_h) * 100, 0),
            "text_preview": _text_of(shape)[:40],
            "is_target": id(shape) in target_ids,
        })
        # Only assign zones from shapes that are injection targets
        if id(shape) not in target_ids:
            continue
        if role == "title" and title_shape is None:
            title_shape = shape
        elif role in ("body", "info"):
            body_shapes.append(shape)

    # Fallback: if classifier didn't find title/body among targets,
    # pick from targets by font size (title) and area (body)
    if not title_shape and targets:
        top_half = [s for s in targets
                    if s.has_text_frame and _shape_top_frac(s, slide_h) < 0.45]
        if top_half:
            title_shape = max(top_half, key=lambda s: _max_font_pt(s))

    if not body_shapes:
        for s in sorted(targets, key=lambda s: _shape_area(s), reverse=True):
            if s != title_shape and s.has_text_frame:
                body_shapes.append(s)
            if len(body_shapes) >= th.body_max_zones:
                break

    # Step 4: Inject content into zones (shapes already cleared by step 2)
    # --- Title ---
    if content_data.title and title_shape:
        _inject_text_simple(title_shape, content_data.title)
        diag["injected_title"] = content_data.title[:50]

    # --- Body ---
    if content_data.body_paragraphs and body_shapes:
        if len(body_shapes) == 1:
            _inject_structured_text(
                body_shapes[0], content_data.body_paragraphs, th, slide_w, slide_h,
            )
            wc = sum(_word_count(p.text) for p in content_data.body_paragraphs)
            diag["injected_body"] = f"{wc} words -> 1 zone"
        else:
            per_zone = max(1, len(content_data.body_paragraphs) // len(body_shapes))
            idx = 0
            for i, zone in enumerate(body_shapes):
                chunk = (content_data.body_paragraphs[idx:]
                         if i == len(body_shapes) - 1
                         else content_data.body_paragraphs[idx:idx + per_zone])
                idx += per_zone
                if chunk:
                    _inject_structured_text(zone, chunk, th, slide_w, slide_h)
            wc = sum(_word_count(p.text) for p in content_data.body_paragraphs)
            diag["injected_body"] = f"{wc} words -> {len(body_shapes)} zones"

    return diag


# ============================================================================
# E. TABLE, CHART & IMAGE HANDLING
# ============================================================================

def _inject_table_cell_text(cell, text: str) -> None:
    """Fill a table cell preserving its formatting."""
    ns_a = _NSMAP["a"]
    tf = cell.text_frame
    if not tf.paragraphs:
        cell.text = text
        return
    # Save format from first paragraph
    first_p = tf.paragraphs[0]._p
    rPr = None
    for r in first_p.findall(f'{{{ns_a}}}r'):
        rp = r.find(f'{{{ns_a}}}rPr')
        if rp is not None:
            rPr = deepcopy(rp)
            break
    pPr = first_p.find(f'{{{ns_a}}}pPr')
    pPr = deepcopy(pPr) if pPr is not None else None

    # Clear and refill
    txBody = tf._txBody
    for p in list(txBody.findall(f'{{{ns_a}}}p')):
        txBody.remove(p)
    new_p = etree.SubElement(txBody, f'{{{ns_a}}}p')
    if pPr is not None:
        new_p.append(pPr)
    new_r = etree.SubElement(new_p, f'{{{ns_a}}}r')
    if rPr is not None:
        new_r.append(rPr)
    new_t = etree.SubElement(new_r, f'{{{ns_a}}}t')
    new_t.text = text


def _add_table_rows(table, count: int) -> None:
    """Clone the last row of a table to add more rows."""
    ns_a = _NSMAP["a"]
    tbl_el = table._tbl
    rows = tbl_el.findall(f'{{{ns_a}}}tr')
    if not rows:
        return
    last_row = rows[-1]
    for _ in range(count):
        new_row = deepcopy(last_row)
        # Clear text in cloned cells
        for tc in new_row.findall(f'{{{ns_a}}}tc'):
            for p in tc.findall(f'.//{{{ns_a}}}p'):
                for r in p.findall(f'{{{ns_a}}}r'):
                    t = r.find(f'{{{ns_a}}}t')
                    if t is not None:
                        t.text = ""
        tbl_el.append(new_row)


def _handle_tables(cloned_slide, content_data: ContentData, slide_w: int, slide_h: int) -> None:
    if not content_data.tables:
        return

    template_tables = [s for s in cloned_slide.shapes if _is_table(s)]
    ct = content_data.tables[0]

    if template_tables:
        tmpl = template_tables[0].table
        c_data = ct["data"]
        if not c_data:
            return

        c_rows, c_cols = len(c_data), len(c_data[0]) if c_data else 0
        t_rows, t_cols = len(tmpl.rows), len(tmpl.columns)

        # Expand rows if needed (up to 2x)
        if c_rows > t_rows and c_rows <= t_rows * 2:
            _add_table_rows(tmpl, c_rows - t_rows)
            t_rows = c_rows

        # Fill cells preserving formatting
        for ri in range(min(c_rows, t_rows)):
            for ci in range(min(c_cols, t_cols)):
                try:
                    _inject_table_cell_text(tmpl.cell(ri, ci), c_data[ri][ci])
                except Exception:
                    pass
        # Clear extra template cells
        for ri in range(c_rows, t_rows):
            for ci in range(t_cols):
                try:
                    _inject_table_cell_text(tmpl.cell(ri, ci), "")
                except Exception:
                    pass
    else:
        try:
            spTree = cloned_slide.shapes._spTree
            spTree.append(deepcopy(ct["element"]))
        except Exception:
            pass


def _handle_charts(
    cloned_slide, content_data: ContentData,
    src_slide, dst_prs: Presentation,
) -> None:
    """Best-effort chart transfer from content slide."""
    if not content_data.charts:
        return
    for chart_info in content_data.charts:
        try:
            chart_part = chart_info["chart_part"]
            # Copy chart part into destination
            chart_rel_type = (
                "http://schemas.openxmlformats.org/officeDocument/2006/relationships/chart"
            )
            new_rid = cloned_slide.part.rels.get_or_add(chart_rel_type, chart_part)
            # Add the chart graphicFrame element
            el = deepcopy(chart_info["element"])
            # Update rId in the element
            for node in el.iter():
                for attr in list(node.attrib.keys()):
                    if node.attrib[attr] in ("rId1", "rId2", "rId3"):
                        node.attrib[attr] = new_rid
            cloned_slide.shapes._spTree.append(el)
        except Exception as exc:
            log.warning("  Chart transfer failed: %s", exc)


def _handle_images(
    cloned_slide, content_data: ContentData,
    slide_w: int, slide_h: int,
) -> None:
    if not content_data.images:
        return
    occupied = []
    for shape in cloned_slide.shapes:
        occupied.append((shape.top or 0) + (shape.height or 0))

    for blob, orig_w, orig_h, _, _ in content_data.images:
        max_bottom = max(occupied, default=int(slide_h * 0.3))
        avail_top = min(max_bottom + int(Pt(10).emu), int(slide_h * 0.85))
        avail_h = slide_h - avail_top
        if avail_h < int(slide_h * 0.1):
            continue

        tw = min(orig_w, int(slide_w * 0.6))
        th_ = min(orig_h, avail_h)
        if orig_w > 0 and orig_h > 0:
            scale = min(tw / orig_w, th_ / orig_h)
            tw, th_ = int(orig_w * scale), int(orig_h * scale)

        try:
            cloned_slide.shapes.add_picture(
                io.BytesIO(blob), (slide_w - tw) // 2, avail_top, tw, th_,
            )
        except Exception:
            pass


# ============================================================================
# F. POST-PROCESSING
# ============================================================================

def _post_process(output_prs: Presentation) -> None:
    sw = output_prs.slide_width or 1
    sh = output_prs.slide_height or 1

    for slide_idx, slide in enumerate(output_prs.slides):
        slide_num = slide_idx + 1

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text

            # Page numbers — text pattern
            if _PAGE_NUM_PATTERN.search(text):
                new = _PAGE_NUM_PATTERN.sub(f"Page {slide_num:02d}", text)
                if new != text:
                    _inject_text_simple(shape, new)
                continue

            # XML-level slide number placeholder
            if _has_placeholder_type(shape, {_PH_SLIDE_NUMBER}):
                _inject_text_simple(shape, str(slide_num))
                continue

            # Date placeholder
            if _has_placeholder_type(shape, {_PH_DATE}):
                _inject_text_simple(shape, date.today().strftime("%Y-%m-%d"))
                continue

            # Dates in footer area
            if _shape_bottom_frac(shape, sh) >= 0.90:
                m = _DATE_PATTERN.search(text)
                if m:
                    new = _DATE_PATTERN.sub(date.today().strftime("%Y-%m-%d"), text)
                    if new != text:
                        _inject_text_simple(shape, new)


def _cleanup_broken_rels(output_prs: Presentation) -> int:
    """Remove broken relationship references that prevent LibreOffice from opening.

    Returns count of removed relationships.
    """
    removed = 0
    for slide in output_prs.slides:
        part = slide.part
        bad_keys: list[str] = []
        for rel_key, rel in part.rels.items():
            try:
                if not rel.is_external:
                    _ = rel.target_part  # will throw if broken
            except Exception:
                bad_keys.append(rel_key)
        for key in bad_keys:
            try:
                del part.rels[key]
                removed += 1
            except Exception:
                pass
    return removed


def _transfer_notes(src_content: ContentData, dst_slide) -> None:
    """Copy speaker notes from content data to the output slide."""
    if not src_content.notes:
        return
    try:
        notes_slide = dst_slide.notes_slide
        tf = notes_slide.notes_text_frame
        if tf:
            existing = tf.text.strip()
            if existing:
                tf.text = existing + "\n\n---\n\n" + src_content.notes
            else:
                tf.text = src_content.notes
    except Exception:
        pass


def _validate_output(output_prs: Presentation) -> list[str]:
    """Validate the output presentation. Returns list of warnings."""
    warnings = []
    for i, slide in enumerate(output_prs.slides):
        shapes = list(slide.shapes)
        if not shapes:
            warnings.append(f"Slide {i+1}: no shapes")
    return warnings


# ============================================================================
# G. DIAGNOSTICS & REPORTING
# ============================================================================

def _print_slide_diagnostic(
    slide_idx: int, total: int, content_data: ContentData,
    template_idx: int, template_struct: str,
    match_score: float, injection_diag: dict,
) -> None:
    print(f"\nSlide {slide_idx+1}/{total}:")
    print(f"  Content type: {content_data.slide_type} "
          f"({content_data.word_count} words, "
          f"{len(content_data.tables)} table(s), "
          f"{len(content_data.images)} image(s))")
    print(f"  Template match: slide {template_idx+1} "
          f"(score={match_score:.0f}, type={template_struct})")

    if injection_diag.get("shapes"):
        print("  Shape classifications:")
        for s in injection_diag["shapes"]:
            preview = f' "{s["text_preview"]}"' if s["text_preview"] else ""
            print(f'    Shape "{s["name"]}" ({s["area_pct"]}% area, '
                  f'top {s["top_pct"]:.0f}%, '
                  f'conf={s["confidence"]}){preview} -> {s["role"]}')

    if injection_diag.get("injected_title"):
        print(f'  Injected: title="{injection_diag["injected_title"]}"')
    if injection_diag.get("injected_body"):
        print(f'  Injected: body ({injection_diag["injected_body"]})')
    print(f'  Cleared: {injection_diag.get("cleared_count", 0)} template text shapes')
    print(f'  Protected: {injection_diag.get("protected_count", 0)} shapes untouched')


# ============================================================================
# INPUT VALIDATION
# ============================================================================

def _validate_input(path: Path, label: str) -> None:
    if not path.exists():
        print(f"Error: {label} not found: {path}", file=sys.stderr)
        sys.exit(1)
    if not path.suffix.lower() == ".pptx":
        print(f"Error: {label} must be a .pptx file: {path}", file=sys.stderr)
        sys.exit(1)
    try:
        with zipfile.ZipFile(str(path), "r") as zf:
            if "[Content_Types].xml" not in zf.namelist():
                print(f"Error: {label} is not a valid PPTX: {path}", file=sys.stderr)
                sys.exit(1)
    except zipfile.BadZipFile:
        print(f"Error: {label} is corrupt or not a ZIP archive: {path}", file=sys.stderr)
        sys.exit(1)

    try:
        prs = Presentation(str(path))
        if len(prs.slides) == 0:
            print(f"Warning: {label} has 0 slides: {path}", file=sys.stderr)
    except Exception as exc:
        print(f"Error: Cannot load {label}: {exc}", file=sys.stderr)
        sys.exit(1)


# ============================================================================
# DESIGN MODE ORCHESTRATOR
# ============================================================================

def apply_design(
    template_path: Path, content_path: Path, output_path: Path,
    config: TransferConfig,
) -> dict[str, Any]:
    """Design mode: clone template slides, inject content. Returns report dict."""
    report: dict[str, Any] = {"mode": "design", "slides": [], "warnings": [], "errors": []}

    print("\n[design] Loading presentations...")
    template_prs = Presentation(str(template_path))
    content_prs = Presentation(str(content_path))
    th = config.thresholds

    sw, sh = template_prs.slide_width, template_prs.slide_height
    ct = len(content_prs.slides)
    tt = len(template_prs.slides)
    print(f"  Template: {tt} slides, {Emu(sw).inches:.1f}\" x {Emu(sh).inches:.1f}\"")
    print(f"  Content:  {ct} slides")

    # Step 1: Extract content
    print("\n[design] Extracting content structure...")
    content_data_list = []
    for i, slide in enumerate(content_prs.slides):
        cd = extract_content(slide, i, ct, content_prs.slide_width, content_prs.slide_height, th)
        content_data_list.append(cd)
        log.debug("  Slide %d: type=%s, words=%d, title=%r, paras=%d, tables=%d, images=%d",
                   i + 1, cd.slide_type, cd.word_count, cd.title[:40],
                   len(cd.body_paragraphs), len(cd.tables), len(cd.images))

    # Step 2: Build mapping
    print("\n[design] Mapping content slides to template slides...")
    if config.slide_map:
        mapping = [config.slide_map.get(str(i + 1), 1) - 1 for i in range(ct)]
        print("  Using manual slide mapping")
    else:
        mapping = build_slide_mapping(content_prs, template_prs, content_data_list, th)

    # Template structures for diagnostics
    t_structs = [_classify_template_structure(s, sw, sh, i, tt)
                 for i, s in enumerate(template_prs.slides)]

    # Step 3: Create output
    print("\n[design] Building output presentation...")
    output_prs = Presentation(str(template_path))

    prs_element = output_prs.slides._sldIdLst
    for sldId in list(prs_element):
        rId = sldId.get(f'{{{_NSMAP["r"]}}}id')
        if rId:
            try:
                output_prs.part.drop_rel(rId)
            except Exception:
                pass
        prs_element.remove(sldId)

    # Step 4: Clone and inject — with per-slide error isolation
    print("\n[design] Cloning and injecting content...")
    success_count = 0
    for ci, cd in enumerate(content_data_list):
        ti = mapping[ci]
        slide_report: dict[str, Any] = {
            "index": ci + 1, "content_type": cd.slide_type,
            "template_slide": ti + 1, "template_type": t_structs[ti],
            "title": cd.title[:80], "word_count": cd.word_count,
            "status": "ok",
        }

        try:
            src_slide = template_prs.slides[ti]
            new_slide = _clone_slide(template_prs, src_slide, output_prs)

            diag = inject_content(new_slide, cd, sw, sh, th)

            _handle_tables(new_slide, cd, sw, sh)
            _handle_charts(new_slide, cd, src_slide, output_prs)
            _handle_images(new_slide, cd, sw, sh)

            if config.preserve_notes:
                _transfer_notes(cd, new_slide)

            slide_report["classifications"] = diag.get("shapes", [])
            slide_report["protected_shapes"] = diag.get("protected_count", 0)
            success_count += 1

            # Print progress
            title_preview = cd.title[:50] if cd.title else "(no title)"
            if config.verbose:
                tinfo = {
                    "struct": t_structs[ti],
                    "words": sum(_word_count(_text_of(s)) for s in src_slide.shapes),
                    "has_table": any(_is_table(s) for s in src_slide.shapes),
                    "is_list": t_structs[ti] in ("list", "grid"),
                }
                score = _match_score(
                    cd.slide_type, tinfo["struct"], ci, ti, ct, tt,
                    cd.word_count, tinfo["words"],
                    len(cd.tables) > 0, tinfo["has_table"],
                    len(cd.body_paragraphs), tinfo["is_list"],
                )
                slide_report["match_score"] = round(score, 1)
                _print_slide_diagnostic(ci, ct, cd, ti, t_structs[ti], score, diag)
            else:
                print(f"  Slide {ci+1}/{ct}: [{cd.slide_type}] "
                      f'"{title_preview}" <- template {ti+1} ({t_structs[ti]})')

        except Exception as exc:
            slide_report["status"] = "error"
            slide_report["error"] = str(exc)
            report["errors"].append(f"Slide {ci+1}: {exc}")
            log.error("Slide %d failed: %s\n%s", ci + 1, exc, traceback.format_exc())
            print(f"  Slide {ci+1}/{ct}: ERROR - {exc}")

            # Insert blank template slide as placeholder
            try:
                fallback_layout = output_prs.slide_layouts[0]
                output_prs.slides.add_slide(fallback_layout)
            except Exception:
                pass

        report["slides"].append(slide_report)

    # Step 5: Post-processing
    print("\n[design] Post-processing...")
    _post_process(output_prs)

    # Step 5b: Clean up broken relationships (LibreOffice compatibility)
    removed_rels = _cleanup_broken_rels(output_prs)
    if removed_rels:
        log.info("  Cleaned up %d broken relationship(s)", removed_rels)

    # Step 6: Validate
    warnings = _validate_output(output_prs)
    report["warnings"] = warnings
    for w in warnings:
        log.warning(w)

    # Step 7: Save
    print(f"\n[design] Saving to {output_path}...")
    output_prs.save(str(output_path))
    print(f"[design] Done! {success_count}/{ct} slides created successfully.")
    if report["errors"]:
        print(f"  {len(report['errors'])} slide(s) had errors — see log for details.")

    return report


# ============================================================================
# RECREATE MODE — analyze → extract → rebuild from scratch
# ============================================================================

# ---- Step 1: Template Style Analyzer ----

@dataclass
class TemplateStyle:
    """Visual DNA extracted from a template PPTX."""
    slide_width: int = 0
    slide_height: int = 0
    heading_font: str = "Montserrat"
    body_font: str = "Lato"
    color_primary: str = "2563EB"
    color_secondary: str = "F97316"
    color_text: str = "111827"
    color_muted: str = "475569"
    color_background: str = "F7F8FB"
    color_card: str = "FFFFFF"
    color_line: str = "D1D5DB"
    logo_blob: bytes | None = None
    logo_content_type: str = "image/png"
    logo_width: int = 0
    logo_height: int = 0
    footer_company: str = ""
    footer_has_confidential: bool = True
    footer_has_page_number: bool = True


def _extract_theme_fonts(prs: Presentation) -> tuple[str, str]:
    """Extract major/minor fonts from theme XML."""
    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
    heading = "Montserrat"
    body = "Lato"
    try:
        master = prs.slide_masters[0]
        theme_el = master.element.find(f".//{{{ns_a}}}theme")
        if theme_el is None:
            # Try reading theme part directly
            for rel in master.part.rels.values():
                if "theme" in str(rel.reltype).lower():
                    theme_xml = rel.target_part.blob
                    theme_el = etree.fromstring(theme_xml)
                    break
        if theme_el is not None:
            major = theme_el.find(f".//{{{ns_a}}}majorFont")
            minor = theme_el.find(f".//{{{ns_a}}}minorFont")
            if major is not None:
                lat = major.find(f"{{{ns_a}}}latin")
                if lat is not None and lat.get("typeface"):
                    heading = lat.get("typeface")
            if minor is not None:
                lat = minor.find(f"{{{ns_a}}}latin")
                if lat is not None and lat.get("typeface"):
                    body = lat.get("typeface")
    except Exception:
        pass

    # Fallback: frequency scan
    if heading == body:
        from collections import Counter
        large_fonts: Counter[str] = Counter()
        body_fonts: Counter[str] = Counter()
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            if run.font.size and run.font.size.pt >= 18:
                                large_fonts[run.font.name] += len(run.text)
                            else:
                                body_fonts[run.font.name] += len(run.text)
        if large_fonts:
            heading = large_fonts.most_common(1)[0][0]
        if body_fonts:
            body = body_fonts.most_common(1)[0][0]
    return heading, body


def _extract_colors(prs: Presentation) -> dict[str, str]:
    """Extract dominant colors from text runs and backgrounds."""
    from collections import Counter
    color_freq: Counter[str] = Counter()
    bg_color = "F7F8FB"

    # Background from first slide
    try:
        fill = prs.slides[0].background.fill
        if fill.type is not None:
            fc = fill.fore_color
            if fc.type is not None and fc.rgb:
                bg_color = str(fc.rgb)
    except Exception:
        pass

    # Text colors
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        c = run.font.color
                        if c and c.type is not None and c.rgb:
                            color_freq[str(c.rgb)] += len(run.text)
                    except (AttributeError, TypeError):
                        pass

    # Classify colors
    text_color = "111827"
    muted_color = "475569"
    primary_color = "2563EB"
    secondary_color = "F97316"

    dark_colors = []
    saturated_accents = []   # True accents (high saturation)
    muted_accents = []       # Grayish mid-tones
    for c, freq in color_freq.most_common(20):
        r, g, b = int(c[:2], 16), int(c[2:4], 16), int(c[4:6], 16)
        brightness = (r + g + b) / 3
        max_ch, min_ch = max(r, g, b), min(r, g, b)
        saturation = (max_ch - min_ch) / max_ch if max_ch > 0 else 0

        if brightness < 80:
            dark_colors.append((c, freq))
        elif c != bg_color and brightness < 240:
            if saturation > 0.4:
                saturated_accents.append((c, freq))
            elif brightness < 160:
                muted_accents.append((c, freq))

    if dark_colors:
        text_color = dark_colors[0][0]
    # Muted = most common grayish mid-tone or second dark
    if muted_accents:
        muted_color = muted_accents[0][0]
    elif len(dark_colors) >= 2:
        muted_color = dark_colors[1][0]
    # Primary = most common saturated accent
    if saturated_accents:
        primary_color = saturated_accents[0][0]
        if len(saturated_accents) >= 2:
            secondary_color = saturated_accents[1][0]

    return {
        "text": text_color, "muted": muted_color,
        "primary": primary_color, "secondary": secondary_color,
        "background": bg_color, "card": "FFFFFF", "line": "D1D5DB",
    }


def _extract_logo(prs: Presentation) -> tuple[bytes | None, str, int, int]:
    """Find the most common image across slides (likely a logo)."""
    from collections import defaultdict
    img_map: dict[int, list] = defaultdict(list)
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or shape.shape_type == 13:
                    blob = shape.image.blob
                    key = hash(blob[:200])
                    img_map[key].append((blob, shape.image.content_type,
                                         shape.width, shape.height))
            except Exception:
                pass

    best_blob = None
    best_ct = "image/png"
    best_w = best_h = 0
    best_count = 0
    for key, occurrences in img_map.items():
        if len(occurrences) > best_count:
            best_count = len(occurrences)
            b, ct, w, h = occurrences[0]
            best_blob, best_ct, best_w, best_h = b, ct, w, h

    if best_count < 2:
        return None, "image/png", 0, 0
    return best_blob, best_ct, best_w, best_h


def _extract_footer_text(prs: Presentation) -> str:
    """Find the most common footer company text."""
    from collections import Counter
    footer_texts: Counter[str] = Counter()
    sh = prs.slide_height
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            bottom = (shape.top or 0) + (shape.height or 0)
            if bottom / sh > 0.90:
                t = shape.text_frame.text.strip()
                if t and not _FOOTER_PATTERNS.match(t) and not re.match(r"^Page\s*\d+$", t, re.I):
                    footer_texts[t] += 1
    if footer_texts:
        return footer_texts.most_common(1)[0][0]
    return ""


def analyze_template(template_path: Path) -> TemplateStyle:
    """Analyze a template PPTX and extract its visual DNA."""
    prs = Presentation(str(template_path))
    style = TemplateStyle()
    style.slide_width = prs.slide_width
    style.slide_height = prs.slide_height

    # Fonts
    style.heading_font, style.body_font = _extract_theme_fonts(prs)

    # Colors
    colors = _extract_colors(prs)
    style.color_primary = colors["primary"]
    style.color_secondary = colors["secondary"]
    style.color_text = colors["text"]
    style.color_muted = colors["muted"]
    style.color_background = colors["background"]
    style.color_card = colors["card"]
    style.color_line = colors["line"]

    # Logo
    blob, ct, w, h = _extract_logo(prs)
    style.logo_blob = blob
    style.logo_content_type = ct
    style.logo_width = w
    style.logo_height = h

    # Footer
    style.footer_company = _extract_footer_text(prs)
    style.footer_has_confidential = True
    style.footer_has_page_number = True

    return style


# ---- Step 2: Content Extractor (reuses existing extract_content) ----

def extract_all_content(
    content_path: Path, th: Thresholds,
) -> list[ContentData]:
    """Extract structured content from every slide in a PPTX."""
    prs = Presentation(str(content_path))
    sw, sh = prs.slide_width, prs.slide_height
    ct = len(prs.slides)
    result = []
    for i, slide in enumerate(prs.slides):
        cd = extract_content(slide, i, ct, sw, sh, th)
        result.append(cd)
    return result


# ---- Step 3: Slide Builder ----

def _rgb(hex_str: str) -> RGBColor:
    """Convert hex string to RGBColor."""
    return RGBColor.from_string(hex_str)


def _style_runs(
    paragraph, *,
    font_name: str, font_size_pt: float,
    bold: bool = False, italic: bool = False,
    color_hex: str = "111827",
) -> None:
    """Apply font properties to every run in a paragraph.

    python-pptx paragraph-level font is a convenience proxy — it doesn't
    always persist into the XML.  Setting on each run is authoritative.
    """
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(font_size_pt)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color_hex)


def _add_background(slide, style: TemplateStyle) -> None:
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(style.color_background)


def _add_decorative_shapes(slide, style: TemplateStyle) -> None:
    """Add corner decorative shapes matching the template style."""
    sw, sh = style.slide_width, style.slide_height

    try:
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MSO_SHAPE

    # Bottom-right ellipse (large, subtle)
    ellipse_size = int(sw * 0.16)
    ellipse = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        sw - int(ellipse_size * 0.7),
        sh - int(ellipse_size * 0.7),
        ellipse_size, ellipse_size,
    )
    fill = ellipse.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(style.color_primary)
    # Set transparency via XML
    ns_a = _NSMAP["a"]
    solid_fill = ellipse._element.find(f".//{{{ns_a}}}solidFill")
    if solid_fill is not None:
        color_el = solid_fill[0] if len(solid_fill) else None
        if color_el is not None:
            alpha = etree.SubElement(color_el, f"{{{ns_a}}}alpha")
            alpha.set("val", "25000")  # 25% opacity
    ellipse.line.fill.background()  # no border

    # Top-right triangle
    tri_w = int(sw * 0.17)
    tri_h = int(sh * 0.33)
    triangle = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        sw - tri_w, 0,
        tri_w, tri_h,
    )
    fill = triangle.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(style.color_line)
    solid_fill = triangle._element.find(f".//{{{ns_a}}}solidFill")
    if solid_fill is not None:
        color_el = solid_fill[0] if len(solid_fill) else None
        if color_el is not None:
            alpha = etree.SubElement(color_el, f"{{{ns_a}}}alpha")
            alpha.set("val", "20000")  # 20% opacity
    triangle.line.fill.background()
    # Flip horizontal so hypotenuse faces left
    triangle.rotation = 180.0


def _add_header(
    slide, style: TemplateStyle, section_label: str,
) -> None:
    """Add accent line and section label above the title area."""
    sw = style.slide_width
    left = int(sw * 0.054)
    # Blue accent line
    line_w = int(sw * 0.038)
    line_h = Pt(3)
    line_top = int(style.slide_height * 0.075)
    line_shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, line_top, line_w, int(line_h),
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = _rgb(style.color_primary)
    line_shape.line.fill.background()

    # Section label (ALL-CAPS, small, primary color)
    if section_label:
        # Sanitize non-ASCII for broad viewer compat
        safe_label = section_label.upper().encode("ascii", "replace").decode("ascii")
        lbl = slide.shapes.add_textbox(
            left, line_top - Pt(16), int(sw * 0.6), Pt(14),
        )
        tf = lbl.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = safe_label
        _style_runs(p, font_name=style.heading_font, font_size_pt=8,
                     bold=True, color_hex=style.color_primary)


def _add_logo(slide, style: TemplateStyle) -> None:
    """Add logo image in the top-left area."""
    if not style.logo_blob:
        return
    try:
        left = int(style.slide_width * 0.024)
        top = int(style.slide_height * 0.030)
        # Scale logo to reasonable size
        max_w = int(style.slide_width * 0.12)
        max_h = int(style.slide_height * 0.05)
        w, h = style.logo_width, style.logo_height
        if w > 0 and h > 0:
            scale = min(max_w / w, max_h / h, 1.0)
            w, h = int(w * scale), int(h * scale)
        else:
            w, h = max_w, max_h
        slide.shapes.add_picture(io.BytesIO(style.logo_blob), left, top, w, h)
    except Exception as exc:
        log.warning("Logo placement failed: %s", exc)


def _add_footer(
    slide, style: TemplateStyle,
    slide_number: int, total_slides: int,
) -> None:
    """Add footer bar with company name, confidential, page number."""
    sw, sh = style.slide_width, style.slide_height
    footer_top = int(sh * 0.94)
    font_size = Pt(7)

    # Company name (left)
    if style.footer_company:
        tb = slide.shapes.add_textbox(
            int(sw * 0.04), footer_top, int(sw * 0.35), Pt(12),
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = style.footer_company
        _style_runs(p, font_name=style.body_font, font_size_pt=7,
                     color_hex=style.color_muted)

    # Confidential (center-right)
    if style.footer_has_confidential:
        tb = slide.shapes.add_textbox(
            int(sw * 0.42), footer_top, int(sw * 0.2), Pt(12),
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = "Confidential"
        _style_runs(p, font_name=style.body_font, font_size_pt=7,
                     color_hex=style.color_muted)

    # Page number (right)
    if style.footer_has_page_number:
        tb = slide.shapes.add_textbox(
            int(sw * 0.90), footer_top, int(sw * 0.07), Pt(12),
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = f"Page {slide_number:02d}"
        p.alignment = PP_ALIGN.RIGHT
        _style_runs(p, font_name=style.body_font, font_size_pt=7,
                     color_hex=style.color_muted)


def _add_title_text(
    slide, style: TemplateStyle, title: str,
    left: int, top: int, width: int,
    font_size_pt: float = 22.0, bold: bool = True,
) -> None:
    """Add title text box."""
    tb = slide.shapes.add_textbox(left, top, width, Pt(font_size_pt + 8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    _style_runs(p, font_name=style.heading_font, font_size_pt=font_size_pt,
                 bold=bold, color_hex=style.color_text)


def _add_body_text(
    slide, style: TemplateStyle,
    paragraphs: list[ParagraphData],
    left: int, top: int, width: int, max_height: int,
) -> None:
    """Add body text box with structured paragraphs."""
    tb = slide.shapes.add_textbox(left, top, width, max_height)
    tf = tb.text_frame
    tf.word_wrap = True

    first = True
    for pd in paragraphs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.text = pd.text
        if pd.bold:
            p.space_before = Pt(10)
            _style_runs(p, font_name=style.heading_font, font_size_pt=14,
                         bold=True, color_hex=style.color_text)
        elif pd.level > 0:
            p.level = pd.level
            _style_runs(p, font_name=style.body_font, font_size_pt=11,
                         color_hex=style.color_muted)
        else:
            _style_runs(p, font_name=style.body_font, font_size_pt=12,
                         color_hex=style.color_text)


def _add_text_blocks(
    slide, style: TemplateStyle,
    text_blocks: list[TextBlock],
) -> None:
    """Recreate positioned text blocks preserving spatial layout from source."""
    sw, sh = style.slide_width, style.slide_height

    for block in text_blocks:
        left = int(sw * block.left_pct / 100)
        top = int(sh * block.top_pct / 100)
        width = int(sw * block.width_pct / 100)
        height = int(sh * block.height_pct / 100)

        # Ensure minimum dimensions
        width = max(width, int(sw * 0.05))
        height = max(height, int(sh * 0.02))

        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True

        first = True
        for pd in block.paragraphs:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()

            p.text = pd.text

            if block.is_label:
                # Small labels: numbers, tags — use accent color, compact
                _style_runs(p, font_name=style.body_font, font_size_pt=10,
                             bold=True, color_hex=style.color_primary)
            elif block.is_heading or pd.bold:
                # Section headings
                p.space_before = Pt(4)
                _style_runs(p, font_name=style.heading_font, font_size_pt=13,
                             bold=True, color_hex=style.color_text)
            elif pd.level > 0:
                p.level = pd.level
                _style_runs(p, font_name=style.body_font, font_size_pt=10,
                             color_hex=style.color_muted)
            else:
                _style_runs(p, font_name=style.body_font, font_size_pt=11,
                             color_hex=style.color_text)


def _add_table(
    slide, style: TemplateStyle,
    table_data: list[list[str]],
    left: int, top: int, width: int, max_height: int,
) -> None:
    """Build a styled table from content data."""
    if not table_data or not table_data[0]:
        return
    rows, cols = len(table_data), len(table_data[0])
    row_height = min(Pt(24), max_height // rows) if rows else Pt(24)

    shape = slide.shapes.add_table(rows, cols, left, top, width, rows * row_height)
    table = shape.table

    # Style each cell
    for ri in range(rows):
        for ci in range(min(cols, len(table_data[ri]) if ri < len(table_data) else 0)):
            cell = table.cell(ri, ci)
            cell.text = table_data[ri][ci] if ri < len(table_data) and ci < len(table_data[ri]) else ""

            # Format
            for p in cell.text_frame.paragraphs:
                _style_runs(p, font_name=style.body_font, font_size_pt=9,
                             color_hex=style.color_text)

            # Header row styling
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(style.color_primary)
                for p in cell.text_frame.paragraphs:
                    _style_runs(p, font_name=style.body_font, font_size_pt=9,
                                 bold=True, color_hex="FFFFFF")
            elif ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(style.color_background)


def _add_content_images(
    slide, style: TemplateStyle,
    images: list[tuple],
    start_top: int,
    has_body_text: bool = True,
) -> None:
    """Place content images on the slide."""
    sw, sh = style.slide_width, style.slide_height
    current_top = start_top

    # If no body text, images can take more space
    max_w_frac = 0.42 if has_body_text else 0.70
    max_h_frac = 0.35 if has_body_text else 0.55

    for img_tuple in images:
        blob = img_tuple[0]
        orig_w, orig_h = img_tuple[1], img_tuple[2]

        max_w = int(sw * max_w_frac)
        max_h = int(sh * max_h_frac)
        avail_h = sh - current_top - int(sh * 0.08)
        if avail_h < int(sh * 0.1):
            break
        max_h = min(max_h, avail_h)

        w, h = orig_w, orig_h
        if w > 0 and h > 0:
            scale = min(max_w / w, max_h / h, 1.0)
            w, h = int(w * scale), int(h * scale)
        else:
            w, h = max_w, max_h

        try:
            left = sw - w - int(sw * 0.04)
            slide.shapes.add_picture(io.BytesIO(blob), left, current_top, w, h)
            current_top += h + Pt(8)
        except Exception:
            pass


def _find_blank_layout(prs: Presentation):
    """Find the blank slide layout (no placeholders, or named 'Blank')."""
    # Prefer layout named "Blank" or "blank"
    for layout in prs.slide_layouts:
        if layout.name.strip().lower() in ("blank", "empty"):
            return layout
    # Fallback: layout with fewest placeholders
    best = prs.slide_layouts[0]
    best_count = len(best.placeholders)
    for layout in prs.slide_layouts:
        if len(layout.placeholders) < best_count:
            best = layout
            best_count = len(layout.placeholders)
    return best


def build_slide(
    prs: Presentation, style: TemplateStyle,
    content: ContentData, slide_number: int, total_slides: int,
) -> None:
    """Build a single output slide from scratch."""
    blank_layout = _find_blank_layout(prs)
    slide = prs.slides.add_slide(blank_layout)
    sw, sh = style.slide_width, style.slide_height

    # Margins
    margin_left = int(sw * 0.054)
    content_width = int(sw * 0.55)
    title_top = int(sh * 0.12)
    body_top = int(sh * 0.22)
    body_max_h = int(sh * 0.65)

    # Background
    _add_background(slide, style)

    # Decorative shapes (skip for title/section slides for cleaner look)
    if content.slide_type not in ("title", "section"):
        _add_decorative_shapes(slide, style)

    # Logo
    _add_logo(slide, style)

    # Section label
    section_label = ""
    if content.slide_type == "title":
        section_label = ""
    elif content.slide_type == "section":
        section_label = ""
    elif content.slide_type == "data":
        section_label = "DATA OVERVIEW"
    elif content.slide_type == "closing":
        section_label = "SUMMARY"
    else:
        # Derive from title
        words = content.title.split()[:3] if content.title else []
        section_label = " ".join(words).upper() if words else "OVERVIEW"

    # --- Layout by slide type ---
    if content.slide_type == "title":
        _build_title_slide(slide, style, content, slide_number, total_slides)
    elif content.slide_type == "section":
        _build_section_slide(slide, style, content, slide_number, total_slides)
    else:
        # Header + footer for content/data/closing
        _add_header(slide, style, section_label)
        _add_footer(slide, style, slide_number, total_slides)

        # Title
        if content.title:
            _add_title_text(
                slide, style, content.title,
                margin_left, title_top, content_width,
            )

        # Determine if we have images to place on the right
        has_images = bool(content.images)

        # Body text — prefer positioned text_blocks for spatial layout
        if content.text_blocks:
            _add_text_blocks(slide, style, content.text_blocks)
        elif content.body_paragraphs:
            bw = int(sw * 0.55) if has_images else int(sw * 0.85)
            _add_body_text(
                slide, style, content.body_paragraphs,
                margin_left, body_top, bw, body_max_h,
            )

        # Tables
        if content.tables:
            table_top = body_top
            if content.body_paragraphs and not content.text_blocks:
                # Estimate body height
                est_lines = sum(1 + len(p.text) // 60 for p in content.body_paragraphs)
                table_top = body_top + int(est_lines * Pt(16))
                table_top = min(table_top, int(sh * 0.55))
            for td in content.tables:
                _add_table(
                    slide, style, td["data"],
                    margin_left, table_top, int(sw * 0.85),
                    int(sh * 0.40),
                )
                break  # Only first table

        # Images — use positioned placement if text_blocks mode, else right-side
        if content.images:
            _add_content_images(slide, style, content.images, body_top,
                                has_body_text=bool(content.body_paragraphs))

    # Speaker notes
    if content.notes:
        try:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            if tf:
                tf.text = content.notes
        except Exception:
            pass


def _build_title_slide(
    slide, style: TemplateStyle, content: ContentData,
    slide_number: int, total_slides: int,
) -> None:
    """Build a title/cover slide."""
    sw, sh = style.slide_width, style.slide_height
    # NOTE: logo already added by build_slide() — no duplicate call here

    # Large centered title
    title_width = int(sw * 0.7)
    title_left = (sw - title_width) // 2
    title_top = int(sh * 0.28)

    if content.title:
        tb = slide.shapes.add_textbox(title_left, title_top, title_width, int(sh * 0.15))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.title
        p.alignment = PP_ALIGN.CENTER
        _style_runs(p, font_name=style.heading_font, font_size_pt=30,
                     bold=True, color_hex=style.color_text)

    # Body content — positioned blocks preserve complex layouts
    if content.text_blocks:
        _add_text_blocks(slide, style, content.text_blocks)
    elif content.body_paragraphs:
        sub_top = title_top + int(sh * 0.18)
        sub_width = int(sw * 0.6)
        sub_left = (sw - sub_width) // 2
        tb = slide.shapes.add_textbox(sub_left, sub_top, sub_width, int(sh * 0.25))
        tf = tb.text_frame
        tf.word_wrap = True
        first = True
        for pd in content.body_paragraphs:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.text = pd.text
            p.alignment = PP_ALIGN.CENTER
            _style_runs(p, font_name=style.body_font, font_size_pt=14,
                         color_hex=style.color_muted)

    # Accent line under title
    line_w = int(sw * 0.08)
    line_left = (sw - line_w) // 2
    line_top = title_top + int(sh * 0.14)
    line_shape = slide.shapes.add_shape(
        1,  # RECTANGLE
        line_left, line_top, line_w, Pt(3),
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = _rgb(style.color_primary)
    line_shape.line.fill.background()

    _add_footer(slide, style, slide_number, total_slides)


def _build_section_slide(
    slide, style: TemplateStyle, content: ContentData,
    slide_number: int, total_slides: int,
) -> None:
    """Build a section divider slide."""
    sw, sh = style.slide_width, style.slide_height
    # NOTE: logo already added by build_slide()

    # Large centered section title
    title_width = int(sw * 0.7)
    title_left = (sw - title_width) // 2
    title_top = int(sh * 0.35)

    if content.title:
        tb = slide.shapes.add_textbox(title_left, title_top, title_width, int(sh * 0.15))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content.title
        p.alignment = PP_ALIGN.CENTER
        _style_runs(p, font_name=style.heading_font, font_size_pt=28,
                     bold=True, color_hex=style.color_text)

    _add_footer(slide, style, slide_number, total_slides)


# ---- Step 4: Orchestrator ----

def apply_recreate(
    template_path: Path, content_path: Path, output_path: Path,
    config: TransferConfig,
) -> dict[str, Any]:
    """Recreate mode: analyze template style, extract content, rebuild from scratch."""
    report: dict[str, Any] = {"mode": "recreate", "slides": [], "warnings": [], "errors": []}
    th = config.thresholds

    # Step 1: Analyze template
    print("\n[recreate] Analyzing template style...")
    style = analyze_template(template_path)
    print(f"  Fonts: heading={style.heading_font}, body={style.body_font}")
    print(f"  Colors: primary=#{style.color_primary}, text=#{style.color_text}, bg=#{style.color_background}")
    print(f"  Logo: {'found' if style.logo_blob else 'not found'} ({style.logo_width}x{style.logo_height})")
    print(f"  Footer: '{style.footer_company}'")

    # Step 2: Extract content
    print("\n[recreate] Extracting content...")
    content_list = extract_all_content(content_path, th)
    ct = len(content_list)
    for i, cd in enumerate(content_list):
        if config.verbose:
            print(f"  Slide {i+1}: type={cd.slide_type}, words={cd.word_count}, "
                  f"title='{cd.title[:40]}', paras={len(cd.body_paragraphs)}, "
                  f"tables={len(cd.tables)}, images={len(cd.images)}")

    # Step 3: Build output — use template as base to preserve theme/masters
    print(f"\n[recreate] Building {ct} slides from scratch...")
    output_prs = Presentation(str(template_path))

    # Remove all template slides but keep masters/theme
    sld_id_lst = output_prs.slides._sldIdLst
    ns_r = _NSMAP["r"]
    for sld_id in list(sld_id_lst):
        r_id = sld_id.get(f"{{{ns_r}}}id")
        if r_id:
            try:
                output_prs.part.drop_rel(r_id)
            except Exception:
                pass
        sld_id_lst.remove(sld_id)

    for i, cd in enumerate(content_list):
        slide_report: dict[str, Any] = {
            "index": i + 1, "content_type": cd.slide_type,
            "title": cd.title[:80] if cd.title else "", "word_count": cd.word_count,
            "status": "ok",
        }
        try:
            build_slide(output_prs, style, cd, i + 1, ct)
            print(f"  Slide {i+1}/{ct}: [{cd.slide_type}] "
                  f'"{cd.title[:50] if cd.title else "(no title)"}"')
        except Exception as exc:
            slide_report["status"] = "error"
            slide_report["error"] = str(exc)
            report["errors"].append(f"Slide {i+1}: {exc}")
            log.error("Slide %d failed: %s\n%s", i + 1, exc, traceback.format_exc())
            print(f"  Slide {i+1}/{ct}: ERROR - {exc}")
            # Add blank slide as fallback
            try:
                output_prs.slides.add_slide(_find_blank_layout(output_prs))
            except Exception:
                pass
        report["slides"].append(slide_report)

    # Save
    print(f"\n[recreate] Saving to {output_path}...")
    output_prs.save(str(output_path))
    success = sum(1 for s in report["slides"] if s["status"] == "ok")
    print(f"[recreate] Done! {success}/{ct} slides created successfully.")

    if config.report_path:
        clean = json.loads(json.dumps(report, default=str))
        config.report_path.write_text(json.dumps(clean, indent=2))
        print(f"Report written to {config.report_path}")

    return report


# ============================================================================
# LAYOUT MODE (backward-compatible)
# ============================================================================

def apply_layout(
    template_path: Path, content_path: Path, output_path: Path,
    config: TransferConfig,
) -> dict[str, Any]:
    print("\n[layout] Loading presentations...")
    template_prs = Presentation(str(template_path))
    print(f"  Template layouts: {[l.name for l in template_prs.slide_layouts]}")
    print("  [layout] Using design-mode pipeline (python-pptx layout limitation)")
    return apply_design(template_path, content_path, output_path, config)


# ============================================================================
# AUTO-DETECTION
# ============================================================================

def detect_mode(template_path: Path) -> str:
    return "recreate"  # New default: always use recreate mode


# ============================================================================
# PUBLIC API
# ============================================================================

def transfer(
    template: Path, content: Path,
    output: Path | None = None,
    config: TransferConfig | None = None,
) -> Presentation | dict:
    """Programmatic API for template transfer.

    Returns the report dict. If output is None, returns (report, Presentation).
    """
    if config is None:
        config = TransferConfig()

    mode = config.mode or detect_mode(template)
    if output is None:
        output = Path("output.pptx")

    if mode == "recreate":
        return apply_recreate(template, content, output, config)
    if mode == "design" or mode == "clone":
        return apply_design(template, content, output, config)
    return apply_layout(template, content, output, config)


# ============================================================================
# CLI: ANALYSIS MODES
# ============================================================================

def _cli_analyze(pptx_path: Path) -> None:
    """Run shape classification on every slide and print results."""
    prs = Presentation(str(pptx_path))
    sw, sh = prs.slide_width, prs.slide_height
    th = Thresholds()

    print(f"\nAnalyzing: {pptx_path}")
    print(f"Slides: {len(prs.slides)}, Size: {Emu(sw).inches:.1f}\" x {Emu(sh).inches:.1f}\"\n")

    for i, slide in enumerate(prs.slides):
        zones = get_slide_zones(slide, sw, sh, th)
        classifications = classify_all_shapes(slide, sw, sh, th)
        total = len(list(slide.shapes))

        print(f"Slide {i+1}/{len(prs.slides)}: {total} shapes")
        print(f"  Zones: title={len(zones['title'])}, body={len(zones['body'])}, "
              f"footer={len(zones['footer'])}, decorative={len(zones['decorative'])}")

        for shape, role, conf in classifications:
            text = _text_of(shape)
            preview = f' "{text[:50]}"' if text else ""
            area = round(_shape_area_pct(shape, sw, sh), 1)
            top = round(_shape_top_frac(shape, sh) * 100)
            print(f'    [{role:11s} {conf:.2f}] "{shape.name}" '
                  f'({area}% area, top {top}%){preview}')
        print()


def _cli_extract(pptx_path: Path) -> None:
    """Extract structured content from every slide and print as JSON."""
    prs = Presentation(str(pptx_path))
    sw, sh = prs.slide_width, prs.slide_height
    th = Thresholds()
    ct = len(prs.slides)

    result = []
    for i, slide in enumerate(prs.slides):
        cd = extract_content(slide, i, ct, sw, sh, th)
        slide_data = {
            "slide": i + 1,
            "slide_type": cd.slide_type,
            "title": cd.title,
            "word_count": cd.word_count,
            "body_paragraphs": [
                {"text": p.text, "bold": p.bold, "level": p.level}
                for p in cd.body_paragraphs
            ],
            "tables": [
                t["data"] for t in cd.tables
            ],
            "images": [
                {"width": img[1], "height": img[2], "blob_size": len(img[0])}
                for img in cd.images
            ],
            "has_chart": cd.has_chart,
            "notes": cd.notes if cd.notes else None,
        }
        result.append(slide_data)

    print(json.dumps(result, indent=2, ensure_ascii=False))


# ============================================================================
# CLI: MAIN
# ============================================================================

def main() -> None:
    parser = argparse.ArgumentParser(
        description="PPTX Template Transfer - apply one deck's visual design to another's content.",
    )
    parser.add_argument("template_pptx", type=Path, help="Template PPTX (or input for --analyze/--extract)")
    parser.add_argument("content_pptx", type=Path, nargs="?", default=None,
                        help="Content PPTX (not needed for --analyze/--extract)")
    parser.add_argument("output_pptx", type=Path, nargs="?", default=None,
                        help="Output PPTX path (not needed for --analyze/--extract)")
    parser.add_argument("--mode", choices=["recreate", "design", "clone", "layout"], default=None)
    parser.add_argument("--slide-map", type=Path, default=None,
                        help='JSON: {"1": 3, "2": 1, ...}')
    parser.add_argument("--layout-map", type=Path, default=None)
    parser.add_argument("--verbose", "-v", action="store_true")
    parser.add_argument("--report", type=Path, default=None,
                        help="Write JSON diagnostics report to this path")
    parser.add_argument("--no-notes", action="store_true",
                        help="Skip speaker notes transfer")
    parser.add_argument("--analyze", action="store_true",
                        help="Analyze a single PPTX: classify every shape on every slide")
    parser.add_argument("--extract", action="store_true",
                        help="Extract structured content from a single PPTX as JSON")

    args = parser.parse_args()

    # Ensure stdout handles unicode
    if hasattr(sys.stdout, "buffer"):
        sys.stdout = io.TextIOWrapper(
            sys.stdout.buffer, encoding="utf-8", errors="replace",
        )

    # Logging
    level = logging.DEBUG if args.verbose else logging.INFO
    logging.basicConfig(level=level, format="%(message)s", stream=sys.stdout)

    # --- Single-file analysis modes ---
    if args.analyze:
        _validate_input(args.template_pptx, "Input")
        _cli_analyze(args.template_pptx)
        return

    if args.extract:
        _validate_input(args.template_pptx, "Input")
        _cli_extract(args.template_pptx)
        return

    # --- Transfer mode: require content + output ---
    if not args.content_pptx or not args.output_pptx:
        parser.error("content_pptx and output_pptx are required for transfer mode")

    _validate_input(args.template_pptx, "Template")
    _validate_input(args.content_pptx, "Content")

    slide_map = None
    if args.slide_map and args.slide_map.exists():
        slide_map = json.loads(args.slide_map.read_text())

    config = TransferConfig(
        mode=args.mode, verbose=args.verbose, slide_map=slide_map,
        preserve_notes=not args.no_notes, report_path=args.report,
    )

    mode = config.mode or detect_mode(args.template_pptx)
    if config.mode is None:
        print(f"Auto-detected mode: {mode}")
    config = TransferConfig(
        mode=mode, verbose=config.verbose, slide_map=config.slide_map,
        preserve_notes=config.preserve_notes, thresholds=config.thresholds,
        report_path=config.report_path,
    )

    if mode == "recreate":
        report = apply_recreate(args.template_pptx, args.content_pptx, args.output_pptx, config)
    elif mode in ("design", "clone"):
        report = apply_design(args.template_pptx, args.content_pptx, args.output_pptx, config)
    else:
        report = apply_layout(args.template_pptx, args.content_pptx, args.output_pptx, config)

    if args.report:
        clean_report = json.loads(json.dumps(report, default=str))
        args.report.write_text(json.dumps(clean_report, indent=2))
        print(f"Report written to {args.report}")


if __name__ == "__main__":
    main()
